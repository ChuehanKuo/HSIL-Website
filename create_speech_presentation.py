#!/usr/bin/env python3
"""
HSIL Hackathon 2026 — Taiwan Hub
Light theme, information-rich, with flowcharts
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# ── Colors (LIGHT PALETTE) ─────────────────────────────
CRIMSON      = RGBColor(0xA8, 0x1D, 0x35)
TEAL         = RGBColor(0x0D, 0x9B, 0x76)
GOLD         = RGBColor(0xD4, 0x8B, 0x0C)
DARK         = RGBColor(0x1E, 0x29, 0x3B)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xFA, 0xFA, 0xFC)
LIGHT_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY     = RGBColor(0xCB, 0xD5, 0xE1)
TEXT_GRAY    = RGBColor(0x64, 0x74, 0x8B)
SOFT_TEAL    = RGBColor(0xE6, 0xFA, 0xF5)
SOFT_CRIM    = RGBColor(0xFD, 0xF0, 0xF3)
SOFT_GOLD    = RGBColor(0xFE, 0xF3, 0xE2)
SOFT_BLUE    = RGBColor(0xEF, 0xF6, 0xFF)
CARD_BG      = RGBColor(0xF8, 0xFA, 0xFC)

ASSETS = "/home/user/HSIL-Website/src/assets"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL = 12

# ── Helpers ─────────────────────────────────────────────
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_shape(slide, shape_type, l, t, w, h, fill=None, line_color=None, line_width=None):
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        if line_width:
            s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s

def rr(slide, l, t, w, h, fill, line_color=None, line_width=None):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill, line_color, line_width)

def rect(slide, l, t, w, h, fill):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, fill)

def oval(slide, l, t, w, h, fill):
    return add_shape(slide, MSO_SHAPE.OVAL, l, t, w, h, fill)

def arrow_right(slide, l, t, w, h, fill):
    return add_shape(slide, MSO_SHAPE.CHEVRON, l, t, w, h, fill)

def tb(slide, l, t, w, h, text, sz=18, clr=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    try: tf.vertical_anchor = anchor
    except: pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = "Calibri"
    p.font.italic = italic
    p.alignment = align
    return box

def multi_text(slide, l, t, w, h, lines):
    """lines = list of (text, size, color, bold, align, italic)"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        txt = item[0]
        sz = item[1] if len(item) > 1 else 16
        clr = item[2] if len(item) > 2 else DARK
        bld = item[3] if len(item) > 3 else False
        al = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        ita = item[5] if len(item) > 5 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = "Calibri"
        p.font.italic = ita
        p.alignment = al
        p.space_before = Pt(0) if i == 0 else Pt(4)
    return box

def slide_num(slide, n):
    tb(slide, Inches(12.3), Inches(7.1), Inches(0.9), Inches(0.3),
       f"{n}/{TOTAL}", sz=10, clr=TEXT_GRAY, align=PP_ALIGN.RIGHT)

def accent_bar(slide, l, t, w, clr=TEAL):
    rect(slide, l, t, w, Pt(4), clr)

def top_bar(slide, clr=CRIMSON):
    rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), clr)

def section_label(slide, l, t, text, clr=CRIMSON):
    tb(slide, l, t, Inches(5), Inches(0.35), text, sz=11, clr=clr, bold=True)

def check_item(slide, x, y, text, sz=14, w=Inches(5.5)):
    """A check mark + text item"""
    tb(slide, x, y, Inches(0.3), Inches(0.35), "✓", sz=sz, clr=TEAL, bold=True)
    tb(slide, x + Inches(0.35), y, w, Inches(0.35), text, sz=sz, clr=DARK)


# ══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)

# Top accent bar
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.1), CRIMSON)
rect(s, Inches(0), Inches(0.1), Inches(13.333), Inches(0.04), TEAL)

# Left decorative band
rect(s, Inches(0), Inches(0.14), Inches(0.5), Inches(7.36), LIGHT_GRAY)

# Logos
try:
    s.shapes.add_picture(os.path.join(ASSETS, "hsil-logo.png"),
                         Inches(1.2), Inches(0.7), height=Inches(0.85))
except: pass

# Institution label
tb(s, Inches(1.2), Inches(1.8), Inches(10), Inches(0.35),
   "HARVARD T.H. CHAN SCHOOL OF PUBLIC HEALTH  ·  HEALTH SYSTEMS INNOVATION LAB",
   sz=11, clr=TEXT_GRAY, bold=True)

# Main title
tb(s, Inches(1.2), Inches(2.5), Inches(10), Inches(1.3),
   "HSIL Hackathon 2026", sz=52, clr=CRIMSON, bold=True)

tb(s, Inches(1.2), Inches(3.7), Inches(10), Inches(0.7),
   "Taiwan Hub  ·  全球健康創新黑客松", sz=28, clr=TEAL, bold=True)

accent_bar(s, Inches(1.2), Inches(4.6), Inches(4.5), GOLD)

tb(s, Inches(1.2), Inches(4.9), Inches(10), Inches(0.5),
   "Building High-Value Health Systems: Leveraging AI", sz=18, clr=GOLD, italic=True)

# Coordinator
tb(s, Inches(1.2), Inches(5.6), Inches(10), Inches(0.4),
   "Chuehan Kuo", sz=20, clr=DARK, bold=True)
tb(s, Inches(1.2), Inches(6.0), Inches(10), Inches(0.35),
   "Founder & Lead Coordinator, Taiwan Hub", sz=14, clr=TEXT_GRAY)

# Bottom info
tb(s, Inches(1.2), Inches(6.6), Inches(5), Inches(0.3),
   "April 10–11, 2026  ·  National Taiwan University  ·  NTU School of Public Health",
   sz=12, clr=TEXT_GRAY)
tb(s, Inches(1.2), Inches(6.9), Inches(5), Inches(0.3),
   "Co-organized by NTU × Harvard HSIL  ·  Free Admission",
   sz=12, clr=TEXT_GRAY)

# Poster on right
try:
    s.shapes.add_picture(os.path.join(ASSETS, "hackathon-poster.png"),
                         Inches(9.5), Inches(1.5), height=Inches(5.0))
except: pass

slide_num(s, 1)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, OFF_WHITE)
top_bar(s, CRIMSON)

section_label(s, Inches(0.8), Inches(0.5), "THE PROBLEM")

tb(s, Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.0),
   "Healthcare Systems Are Under Immense Pressure", sz=36, clr=CRIMSON, bold=True)

accent_bar(s, Inches(0.8), Inches(1.9), Inches(3), CRIMSON)

# Three stat cards
stats = [
    ("5 Million", "people die every year", "from poor-quality healthcare\nin low & middle-income countries",
     SOFT_CRIM, CRIMSON),
    ("$8.6 Trillion", "spent on health globally", "each year — with significant\ninefficiency and waste",
     SOFT_GOLD, GOLD),
    ("3.6 Billion", "people lack access", "to essential health services\nworldwide (WHO estimate)",
     SOFT_TEAL, TEAL),
]

for i, (big_num, subtitle, detail, card_bg, accent) in enumerate(stats):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(2.4)

    card = rr(s, x, y, Inches(3.8), Inches(3.2), card_bg, line_color=MID_GRAY, line_width=0.5)
    accent_bar(s, x + Inches(0.3), y + Inches(0.3), Inches(1.5), accent)

    tb(s, x + Inches(0.3), y + Inches(0.6), Inches(3.2), Inches(0.7),
       big_num, sz=34, clr=accent, bold=True)
    tb(s, x + Inches(0.3), y + Inches(1.3), Inches(3.2), Inches(0.4),
       subtitle, sz=16, clr=DARK, bold=True)

    box = s.shapes.add_textbox(x + Inches(0.3), y + Inches(1.8), Inches(3.2), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(detail.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = "Calibri"

# Bottom insight
tb(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
   "AI has the potential to transform how we deliver, manage, and improve healthcare — but it needs fresh ideas from diverse minds.",
   sz=16, clr=DARK, italic=True)

slide_num(s, 2)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — THE QUESTION
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, TEAL)

# Center the question visually
tb(s, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
   "What if you could help fix it?", sz=48, clr=DARK, bold=True, align=PP_ALIGN.CENTER)

accent_bar(s, Inches(5.5), Inches(3.0), Inches(2.3), TEAL)

tb(s, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6),
   "What if all you needed was an idea, a team, and 2 days?", sz=24, clr=TEAL, align=PP_ALIGN.CENTER)

# Three "no barrier" cards
barriers = [
    ("No Coding Required", "不需要程式背景", SOFT_TEAL, TEAL),
    ("No Medical Degree Needed", "不需要醫學背景", SOFT_CRIM, CRIMSON),
    ("Just Your Creativity", "只需要你的創意和熱情", SOFT_GOLD, GOLD),
]

for i, (en, zh, bg_clr, accent) in enumerate(barriers):
    x = Inches(1.5) + i * Inches(3.6)
    y = Inches(4.5)
    card = rr(s, x, y, Inches(3.2), Inches(1.5), bg_clr, line_color=accent, line_width=1)
    tb(s, x + Inches(0.3), y + Inches(0.3), Inches(2.6), Inches(0.5),
       en, sz=17, clr=accent, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.3), y + Inches(0.85), Inches(2.6), Inches(0.4),
       zh, sz=14, clr=TEXT_GRAY, align=PP_ALIGN.CENTER)

slide_num(s, 3)


# ══════════════════════════════════════════════════════════
# SLIDE 4 — THE ANSWER: HSIL Hackathon Reveal
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, GOLD)

section_label(s, Inches(0.8), Inches(0.5), "THE ANSWER", GOLD)

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.8),
   "HSIL Hackathon — The World's Largest Health Innovation Hackathon",
   sz=32, clr=CRIMSON, bold=True)

accent_bar(s, Inches(0.8), Inches(1.8), Inches(3), GOLD)

# Description paragraph (from website)
tb(s, Inches(0.8), Inches(2.2), Inches(7.5), Inches(2.5),
   "HSIL Hackathon 是由哈佛大學公共衛生學院 Health Systems Innovation Lab 發起的全球健康創新黑客松，"
   "每年在世界各地 30 多個城市同步舉辦，匯集數千名來自不同領域的參與者，共同探索健康系統的未來。\n\n"
   "這不是一般的程式競賽——我們著重跨領域團隊合作，結合醫學、公共衛生、商管、設計、工程、法律、"
   "社會科學等背景的參賽者，在兩天內針對真實的全球健康挑戰，運用數據與 AI 工具提出可落地的創新解方。",
   sz=14, clr=DARK)

# Big stats on the right
big_stats = [
    ("7th", "Edition\n第七屆", CRIMSON),
    ("30+", "Cities Worldwide\n全球城市", TEAL),
    ("50+", "Countries\n參與國家", GOLD),
    ("5,000+", "Participants\n歷年參與者", CRIMSON),
]

for i, (num, label, accent) in enumerate(big_stats):
    x = Inches(8.8)
    y = Inches(2.0) + i * Inches(1.3)

    card = rr(s, x, y, Inches(3.8), Inches(1.1), LIGHT_GRAY)
    tb(s, x + Inches(0.3), y + Inches(0.1), Inches(1.5), Inches(0.9),
       num, sz=30, clr=accent, bold=True, anchor=MSO_ANCHOR.MIDDLE)

    box = s.shapes.add_textbox(x + Inches(1.8), y + Inches(0.15), Inches(1.8), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK if j == 0 else TEXT_GRAY
        p.font.bold = (j == 0)
        p.font.name = "Calibri"

# Quote
rr(s, Inches(0.8), Inches(5.2), Inches(7.5), Inches(1.3), SOFT_CRIM)
rect(s, Inches(0.8), Inches(5.2), Inches(0.1), Inches(1.3), CRIMSON)
tb(s, Inches(1.2), Inches(5.4), Inches(6.8), Inches(0.9),
   "國立臺灣大學 × 哈佛大學聯合主辦\n不需要程式背景，只需要你對健康議題的熱情和創新思維。",
   sz=15, clr=CRIMSON, bold=True)

# Logo
try:
    s.shapes.add_picture(os.path.join(ASSETS, "hsil-logo.png"),
                         Inches(9.0), Inches(6.3), height=Inches(0.6))
except: pass

slide_num(s, 4)


# ══════════════════════════════════════════════════════════
# SLIDE 5 — 2026 THEME + CHALLENGES
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, OFF_WHITE)
top_bar(s, TEAL)

section_label(s, Inches(0.8), Inches(0.5), "2026 THEME", TEAL)

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "運用 AI 打造高價值醫療體系", sz=34, clr=DARK, bold=True)
tb(s, Inches(0.8), Inches(1.65), Inches(11), Inches(0.4),
   "Building High-Value Health Systems: Leveraging AI", sz=16, clr=TEAL, italic=True)

accent_bar(s, Inches(0.8), Inches(2.2), Inches(3), TEAL)

tb(s, Inches(0.8), Inches(2.5), Inches(11), Inches(0.5),
   "今年的黑客松聚焦於如何運用人工智慧 (AI) 技術來改善健康系統。參賽團隊將針對以下八大領域提出創新解方：",
   sz=14, clr=TEXT_GRAY)

# 8 challenge cards with Chinese descriptions
challenges = [
    ("電子病歷 (EHR)\n數據分析", "利用 AI 分析電子病歷\n改善臨床決策支持", SOFT_TEAL, TEAL),
    ("AI 診斷工具", "開發輔助診斷工具\n提升早期疾病偵測能力", SOFT_CRIM, CRIMSON),
    ("臨床對話機器人", "建立智慧健康諮詢系統\n改善醫病溝通", SOFT_GOLD, GOLD),
    ("醫療人力短缺", "運用 AI 優化人力配置\n解決醫護人力不足問題", SOFT_BLUE, RGBColor(0x3B, 0x82, 0xF6)),
    ("心血管疾病", "運用數據驅動方法改善\n心血管疾病預防與管理", SOFT_CRIM, CRIMSON),
    ("癌症照護", "AI 輔助癌症篩檢、診斷\n及個人化治療方案", SOFT_TEAL, TEAL),
    ("心理健康", "利用科技改善心理健康\n服務的可及性與品質", SOFT_GOLD, GOLD),
    ("傳染病防治", "AI 驅動的流行病監測\n與因應策略", SOFT_BLUE, RGBColor(0x3B, 0x82, 0xF6)),
]

for i, (title, desc, bg_clr, accent) in enumerate(challenges):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + col * Inches(3.1)
    y = Inches(3.2) + row * Inches(2.05)

    card = rr(s, x, y, Inches(2.85), Inches(1.8), bg_clr, line_color=accent, line_width=0.75)

    # Title
    box = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(2.45), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(title.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = accent
        p.font.bold = True
        p.font.name = "Calibri"

    # Desc
    box2 = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.85), Inches(2.45), Inches(0.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"

slide_num(s, 5)


# ══════════════════════════════════════════════════════════
# SLIDE 6 — TAIWAN HUB: All Key Details
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, CRIMSON)

section_label(s, Inches(0.8), Inches(0.5), "TAIWAN HUB — 台灣站")

tb(s, Inches(0.8), Inches(1.0), Inches(7), Inches(0.8),
   "Taiwan's First Year — 台灣首次加入", sz=34, clr=CRIMSON, bold=True)

tb(s, Inches(0.8), Inches(1.8), Inches(7), Inches(0.4),
   "台灣首次加入全球 30+ 城市同步舉辦的健康創新盛會", sz=14, clr=TEXT_GRAY)

accent_bar(s, Inches(0.8), Inches(2.3), Inches(2.5), TEAL)

# Quick Facts grid (8 items from website)
facts = [
    ("活動日期", "2026 年 4 月 10–11 日（週五、六）", TEAL),
    ("活動地點", "國立臺灣大學 公共衛生學院", TEAL),
    ("活動時長", "兩天（每天約 8 小時）", CRIMSON),
    ("團隊人數", "3–5 人一組（可個人報名，現場組隊）", CRIMSON),
    ("最終產出", "3 分鐘 Pitch 簡報 + 2 分鐘 Q&A", GOLD),
    ("活動語言", "中文為主、部分資料為英文", GOLD),
    ("參加費用", "完全免費（包含兩天午餐與活動物資）", TEAL),
    ("報名截止", "2026 年 3 月 31 日", CRIMSON),
]

for i, (label, value, accent) in enumerate(facts):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.1)
    y = Inches(2.7) + row * Inches(1.08)

    card = rr(s, x, y, Inches(5.7), Inches(0.85), LIGHT_GRAY)
    # Accent dot
    oval(s, x + Inches(0.25), y + Inches(0.3), Inches(0.22), Inches(0.22), accent)

    tb(s, x + Inches(0.6), y + Inches(0.1), Inches(1.5), Inches(0.35),
       label, sz=12, clr=accent, bold=True)
    tb(s, x + Inches(0.6), y + Inches(0.42), Inches(4.8), Inches(0.35),
       value, sz=14, clr=DARK)

# Contact info at bottom
tb(s, Inches(0.8), Inches(7.0), Inches(10), Inches(0.3),
   "IG: @hsilhackathon.taiwan  ·  Email: hsilhackathon.taiwan@gmail.com  ·  Co-organized by NTU × Harvard HSIL",
   sz=11, clr=TEXT_GRAY)

slide_num(s, 6)


# ══════════════════════════════════════════════════════════
# SLIDE 7 — THE JOURNEY: Pipeline Flowchart
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, OFF_WHITE)
top_bar(s, GOLD)

section_label(s, Inches(0.8), Inches(0.5), "THE JOURNEY — 參賽流程", GOLD)

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "From Registration to Harvard Demo Day", sz=32, clr=DARK, bold=True)

accent_bar(s, Inches(0.8), Inches(1.7), Inches(3), GOLD)

# ── FLOWCHART: 5 chevron-style steps ──
steps_data = [
    ("1", "線上報名", "Now → 3/31", "填寫 Airtable\n報名表單", TEAL, SOFT_TEAL),
    ("2", "錄取通知", "4 月初", "公布名單\n寄送行前指南", RGBColor(0x3B, 0x82, 0xF6), SOFT_BLUE),
    ("3", "兩天黑客松", "4/10–11", "組隊·腦力激盪\n導師指導·方案發想", CRIMSON, SOFT_CRIM),
    ("4", "團隊 Pitch", "4/11 下午", "3 分鐘簡報\n2 分鐘 Q&A", GOLD, SOFT_GOLD),
    ("5", "哈佛 Bootcamp", "5–6 月", "前兩名培訓\n6/19 哈佛 Demo Day", CRIMSON, SOFT_CRIM),
]

step_w = Inches(2.2)
step_h = Inches(3.5)
gap = Inches(0.25)
start_x = Inches(0.6)
base_y = Inches(2.3)

for i, (num, title, date, detail, accent, bg_clr) in enumerate(steps_data):
    x = start_x + i * (step_w + gap)

    # Card
    card = rr(s, x, base_y, step_w, step_h, bg_clr, line_color=accent, line_width=1)

    # Number circle
    circle = oval(s, x + Inches(0.7), base_y + Inches(0.25), Inches(0.7), Inches(0.7), accent)
    tb(s, x + Inches(0.7), base_y + Inches(0.25), Inches(0.7), Inches(0.7),
       num, sz=24, clr=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Title
    tb(s, x + Inches(0.1), base_y + Inches(1.15), Inches(2.0), Inches(0.5),
       title, sz=16, clr=accent, bold=True, align=PP_ALIGN.CENTER)

    # Date
    tb(s, x + Inches(0.1), base_y + Inches(1.6), Inches(2.0), Inches(0.3),
       date, sz=11, clr=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # Detail
    box = s.shapes.add_textbox(x + Inches(0.15), base_y + Inches(2.1), Inches(1.9), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(detail.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    # Arrow between steps
    if i < len(steps_data) - 1:
        ax = x + step_w
        ay = base_y + Inches(1.5)
        tb(s, ax, ay, gap, Inches(0.4), "→", sz=20, clr=TEXT_GRAY, align=PP_ALIGN.CENTER)

# Bottom note
rr(s, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8), SOFT_GOLD, line_color=GOLD, line_width=0.5)
tb(s, Inches(1.1), Inches(6.35), Inches(11), Inches(0.5),
   "兩天的黑客松，有可能帶你走到哈佛。 — Top 2 teams represent Taiwan at Harvard Demo Day (June 19, 2026)",
   sz=15, clr=GOLD, bold=True, align=PP_ALIGN.CENTER)

slide_num(s, 7)


# ══════════════════════════════════════════════════════════
# SLIDE 8 — DAY 1 & DAY 2 DETAILED SCHEDULE
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, CRIMSON)

section_label(s, Inches(0.8), Inches(0.5), "DETAILED SCHEDULE — 活動時程")

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6),
   "Two Days at NTU — April 10 & 11, 2026", sz=30, clr=CRIMSON, bold=True)

accent_bar(s, Inches(0.8), Inches(1.6), Inches(2.5), CRIMSON)

# Day 1 column
day1_x = Inches(0.8)
rr(s, day1_x, Inches(1.9), Inches(5.9), Inches(0.55), CRIMSON)
tb(s, day1_x + Inches(0.3), Inches(1.95), Inches(5.3), Inches(0.45),
   "Day 1 — 4 月 10 日（週五）", sz=16, clr=WHITE, bold=True)

day1_schedule = [
    ("08:00–09:00", "報到與入場"),
    ("09:00–09:30", "開幕式與歡迎致詞"),
    ("09:30–10:00", "問題定義與挑戰說明"),
    ("10:00–10:30", "團隊組成"),
    ("10:30–12:30", "團隊腦力激盪與問題探索"),
    ("12:30–13:30", "午餐與交流"),
    ("13:30–17:00", "方案發想與原型開發"),
    ("17:00–17:30", "導師諮詢時段"),
]

for i, (time, title) in enumerate(day1_schedule):
    y = Inches(2.6) + i * Inches(0.56)
    bg_clr = LIGHT_GRAY if i % 2 == 0 else WHITE
    rr(s, day1_x, y, Inches(5.9), Inches(0.48), bg_clr)
    tb(s, day1_x + Inches(0.2), y + Inches(0.08), Inches(1.6), Inches(0.35),
       time, sz=11, clr=TEXT_GRAY, bold=True)
    tb(s, day1_x + Inches(2.0), y + Inches(0.08), Inches(3.5), Inches(0.35),
       title, sz=12, clr=DARK)

# Day 2 column
day2_x = Inches(7.0)
rr(s, day2_x, Inches(1.9), Inches(5.9), Inches(0.55), TEAL)
tb(s, day2_x + Inches(0.3), Inches(1.95), Inches(5.3), Inches(0.45),
   "Day 2 — 4 月 11 日（週六）", sz=16, clr=WHITE, bold=True)

day2_schedule = [
    ("08:30–09:00", "報到"),
    ("09:00–12:00", "方案精煉與簡報準備"),
    ("12:00–13:00", "午餐"),
    ("13:00–13:30", "簡報練習"),
    ("13:30–15:30", "團隊 Pitch 簡報"),
    ("15:30–16:00", "評審討論"),
    ("16:00–16:30", "頒獎典禮與閉幕式"),
    ("16:30–17:00", "交流與合影"),
]

for i, (time, title) in enumerate(day2_schedule):
    y = Inches(2.6) + i * Inches(0.56)
    bg_clr = LIGHT_GRAY if i % 2 == 0 else WHITE
    rr(s, day2_x, y, Inches(5.9), Inches(0.48), bg_clr)
    tb(s, day2_x + Inches(0.2), y + Inches(0.08), Inches(1.6), Inches(0.35),
       time, sz=11, clr=TEXT_GRAY, bold=True)
    tb(s, day2_x + Inches(2.0), y + Inches(0.08), Inches(3.5), Inches(0.35),
       title, sz=12, clr=DARK)

# Post event
rr(s, Inches(0.8), Inches(7.0), Inches(12.1), Inches(0.35), SOFT_GOLD)
tb(s, Inches(1.1), Inches(7.0), Inches(11.5), Inches(0.35),
   "Post-Event:  5 月 — 優勝隊伍方案精煉  ·  6/19 — 哈佛 Demo Day 向全球評審展示方案",
   sz=12, clr=GOLD, bold=True)

slide_num(s, 8)


# ══════════════════════════════════════════════════════════
# SLIDE 9 — WHO CAN JOIN + ELIGIBILITY
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, OFF_WHITE)
top_bar(s, TEAL)

section_label(s, Inches(0.8), Inches(0.5), "WHO CAN JOIN — 誰可以參加", TEAL)

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "This Hackathon Is For Everyone", sz=34, clr=DARK, bold=True)

tb(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
   "HSIL Hackathon 歡迎所有對健康系統創新有興趣的人參加，不限科系、不限年級、不需程式背景。",
   sz=14, clr=TEXT_GRAY)

accent_bar(s, Inches(0.8), Inches(2.15), Inches(2.5), TEAL)

# Disciplines
disciplines = [
    ("醫學 Medicine", CRIMSON, SOFT_CRIM),
    ("公衛 Public Health", TEAL, SOFT_TEAL),
    ("資工 Computer Science", RGBColor(0x3B, 0x82, 0xF6), SOFT_BLUE),
    ("商管 Business", GOLD, SOFT_GOLD),
    ("設計 Design", RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0xF3, 0xE8, 0xFF)),
    ("工程 Engineering", TEAL, SOFT_TEAL),
    ("法律 Law", CRIMSON, SOFT_CRIM),
    ("社科 Social Science", GOLD, SOFT_GOLD),
]

for i, (name, accent, bg_clr) in enumerate(disciplines):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + col * Inches(3.1)
    y = Inches(2.5) + row * Inches(0.9)

    card = rr(s, x, y, Inches(2.8), Inches(0.7), bg_clr, line_color=accent, line_width=0.75)
    tb(s, x + Inches(0.2), y + Inches(0.12), Inches(2.4), Inches(0.45),
       name, sz=14, clr=accent, bold=True, align=PP_ALIGN.CENTER)

# Eligibility requirements (Chinese from website)
requirements = [
    "大學生、研究生、博士生皆可報名",
    "不限科系 — 醫學、公衛、商管、設計、工程、法律、社會科學⋯都歡迎",
    "不需會寫程式",
    "可個人報名，活動當天現場組隊（3–5 人一組）",
    "須全程參與兩天活動（4/10–11）",
    "須在 3/31 前完成線上報名表單",
]

rr(s, Inches(0.8), Inches(4.5), Inches(12.0), Inches(2.8), WHITE, line_color=MID_GRAY, line_width=0.5)
tb(s, Inches(1.1), Inches(4.6), Inches(5), Inches(0.35),
   "Eligibility Requirements 報名資格", sz=14, clr=TEAL, bold=True)

for i, req in enumerate(requirements):
    y = Inches(5.05) + i * Inches(0.36)
    check_item(s, Inches(1.1), y, req, sz=13, w=Inches(10))

slide_num(s, 9)


# ══════════════════════════════════════════════════════════
# SLIDE 10 — WHY JOIN (6 reasons + FAQ highlights)
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, GOLD)

section_label(s, Inches(0.8), Inches(0.5), "WHY JOIN — 為什麼要參加", GOLD)

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "6 Reasons You Can't Miss This", sz=32, clr=CRIMSON, bold=True)

accent_bar(s, Inches(0.8), Inches(1.65), Inches(2.5), GOLD)

# 6 highlight cards (from website HighlightsSection)
highlights = [
    ("完全免費", "無任何報名費或參加費用\n包含兩天午餐與活動物資", TEAL, SOFT_TEAL),
    ("不需程式背景", "重視創意思維與跨領域合作\n產出是 Pitch 簡報，不需寫 Code", CRIMSON, SOFT_CRIM),
    ("哈佛 Bootcamp 機會", "優勝隊伍將獲邀參加哈佛大學\nDemo Day 及 Bootcamp 培訓計畫", GOLD, SOFT_GOLD),
    ("專業導師指導", "活動期間配有來自醫療、科技、\n商業領域的專業導師全程指導", TEAL, SOFT_TEAL),
    ("跨領域交流", "認識來自醫學、設計、工程、商管、\n法律等不同領域的優秀夥伴", CRIMSON, SOFT_CRIM),
    ("解決真實問題", "針對台灣及全球健康系統面臨的\n真實挑戰提出解方", GOLD, SOFT_GOLD),
]

for i, (title, desc, accent, bg_clr) in enumerate(highlights):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.1) + row * Inches(2.2)

    card = rr(s, x, y, Inches(3.8), Inches(1.9), bg_clr, line_color=accent, line_width=0.75)

    tb(s, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.4),
       title, sz=17, clr=accent, bold=True)

    box = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"

# FAQ highlights at bottom
rr(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.7), LIGHT_GRAY)
tb(s, Inches(1.1), Inches(6.7), Inches(11.2), Inches(0.5),
   "FAQ:  智慧財產權完全歸屬參賽團隊  ·  團隊簡報可選擇中文或英文  ·  建議自備筆電（非強制）",
   sz=12, clr=TEXT_GRAY)

slide_num(s, 10)


# ══════════════════════════════════════════════════════════
# SLIDE 11 — JUDGING CRITERIA
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, OFF_WHITE)
top_bar(s, CRIMSON)

section_label(s, Inches(0.8), Inches(0.5), "JUDGING — 評審標準")

tb(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.7),
   "How Teams Are Evaluated", sz=32, clr=CRIMSON, bold=True)

tb(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.4),
   "6 criteria  ×  5 points each  =  30 points total  ·  總分：30 分", sz=15, clr=TEXT_GRAY)

accent_bar(s, Inches(0.8), Inches(2.1), Inches(2), CRIMSON)

# 6 criteria with full Chinese descriptions from website
criteria = [
    ("創新性", "Innovation", "5", "方案是否具有原創性與創新思維？", CRIMSON),
    ("可行性", "Feasibility", "5", "方案在技術和資源上是否可行？", TEAL),
    ("健康影響力", "Health Impact", "5", "對健康系統的改善潛力有多大？", GOLD),
    ("可擴展性", "Scalability", "5", "方案能否推廣到更大規模？", CRIMSON),
    ("簡報品質", "Presentation", "5", "團隊的表達與展示是否清晰有力？", TEAL),
    ("AI 應用", "AI Integration", "5", "AI 技術是否被有效且合理地應用？", GOLD),
]

for i, (zh, en, pts, desc, accent) in enumerate(criteria):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.5) + row * Inches(2.3)

    card = rr(s, x, y, Inches(3.8), Inches(2.0), WHITE, line_color=accent, line_width=1)

    # Score circle
    circle = oval(s, x + Inches(0.25), y + Inches(0.25), Inches(0.65), Inches(0.65), accent)
    tb(s, x + Inches(0.25), y + Inches(0.25), Inches(0.65), Inches(0.65),
       pts, sz=22, clr=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Title
    tb(s, x + Inches(1.1), y + Inches(0.25), Inches(2.5), Inches(0.4),
       zh, sz=17, clr=accent, bold=True)
    tb(s, x + Inches(1.1), y + Inches(0.6), Inches(2.5), Inches(0.3),
       en, sz=12, clr=TEXT_GRAY, italic=True)

    # Description
    tb(s, x + Inches(0.25), y + Inches(1.15), Inches(3.3), Inches(0.6),
       desc, sz=13, clr=DARK)

slide_num(s, 11)


# ══════════════════════════════════════════════════════════
# SLIDE 12 — CLOSING + CTA
# ══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)

# Top double accent
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.1), CRIMSON)
rect(s, Inches(0), Inches(0.1), Inches(13.333), Inches(0.04), TEAL)

# Big question
tb(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.0),
   "準備好改變世界了嗎？", sz=46, clr=CRIMSON, bold=True, align=PP_ALIGN.CENTER)

tb(s, Inches(0.5), Inches(1.8), Inches(12.3), Inches(0.6),
   "The future of healthcare needs people like you.", sz=22, clr=TEAL, bold=True, align=PP_ALIGN.CENTER)

accent_bar(s, Inches(5.5), Inches(2.7), Inches(2.3), GOLD)

# Key info recap cards
recap = [
    ("Registration Deadline\n報名截止", "March 31, 2026\n2026/3/31", CRIMSON, SOFT_CRIM),
    ("Event Date\n活動日期", "April 10–11, 2026\n4月10-11日", TEAL, SOFT_TEAL),
    ("Venue\n活動地點", "NTU School of\nPublic Health", GOLD, SOFT_GOLD),
    ("Cost\n費用", "Completely\nFREE", TEAL, SOFT_TEAL),
]

for i, (label, value, accent, bg_clr) in enumerate(recap):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(3.1)

    card = rr(s, x, y, Inches(2.85), Inches(1.8), bg_clr, line_color=accent, line_width=1)

    box = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(2.45), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    for j, line in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = accent
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    box2 = s.shapes.add_textbox(x + Inches(0.2), y + Inches(0.85), Inches(2.45), Inches(0.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    for j, line in enumerate(value.split('\n')):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# Contact & links
rr(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.0), LIGHT_GRAY)
tb(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
   "Registration: airtable.com (link on our website)  ·  IG: @hsilhackathon.taiwan  ·  Email: hsilhackathon.taiwan@gmail.com",
   sz=14, clr=DARK, bold=True, align=PP_ALIGN.CENTER)
tb(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.35),
   "名額有限。台灣第一次。不要錯過。", sz=14, clr=CRIMSON, bold=True, align=PP_ALIGN.CENTER)

# Sign-off
tb(s, Inches(0.8), Inches(6.4), Inches(5), Inches(0.35),
   "Chuehan Kuo  ·  Founder & Lead Coordinator, Taiwan Hub", sz=13, clr=TEXT_GRAY)

tb(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
   "Thank you  ·  謝謝", sz=20, clr=DARK, bold=True, align=PP_ALIGN.CENTER)

# Logos at bottom
try:
    s.shapes.add_picture(os.path.join(ASSETS, "hsil-logo.png"),
                         Inches(5.0), Inches(6.3), height=Inches(0.5))
except: pass

slide_num(s, 12)


# ── Save ────────────────────────────────────────────────
output = "/home/user/HSIL-Website/HSIL_Hackathon_2026_Speech.pptx"
prs.save(output)
print(f"Presentation saved: {output}")
print(f"Total slides: {TOTAL}")
