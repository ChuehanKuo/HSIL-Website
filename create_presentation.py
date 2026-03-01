#!/usr/bin/env python3
"""
HSIL Hackathon 2026 — Taiwan Hub
Presentation generator for AI class (5-10 min)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────
CRIMSON      = RGBColor(0xA8, 0x1D, 0x35)  # Harvard crimson
TEAL         = RGBColor(0x0D, 0x9B, 0x76)  # Accent teal
GOLD         = RGBColor(0xD4, 0x8B, 0x0C)  # Accent gold
DARK         = RGBColor(0x1A, 0x20, 0x2E)  # Near-black
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xF1, 0xF5, 0xF9)
MED_GRAY     = RGBColor(0x64, 0x74, 0x8B)
SOFT_TEAL_BG = RGBColor(0xE6, 0xFA, 0xF5)
SOFT_CRIM_BG = RGBColor(0xFD, 0xF0, 0xF3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

ASSETS = "/home/user/HSIL-Website/src/assets"

# ── Helpers ─────────────────────────────────────────────
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        # Set transparency via XML
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_el.set('val', str(int(alpha * 1000)))  # e.g. 50000 for 50%
    return shape

def rounded_rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def text_box(slide, left, top, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_paragraph(text_frame, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def add_bullet_list(slide, left, top, w, h, items, size=18, color=DARK, bullet_color=TEAL):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.level = 0
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        from lxml import etree
        buFont = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
        buFont.set('typeface', 'Arial')
        buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        buChar.set('char', '●')
        buClr = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buClr')
        srgb = etree.SubElement(buClr, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        srgb.set('val', str(bullet_color))
    return txBox

def accent_line(slide, left, top, w, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def slide_number(slide, num, total):
    text_box(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
             f"{num} / {total}", size=11, color=MED_GRAY, align=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 10

# ══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Decorative shapes
rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, CRIMSON)
rect(slide, Inches(0.4), Inches(0), Inches(0.15), SLIDE_H, TEAL)

# Logos
try:
    slide.shapes.add_picture(os.path.join(ASSETS, "hsil-logo.png"),
                             Inches(1.2), Inches(0.6), height=Inches(1.0))
except:
    pass

# Title
text_box(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.2),
         "HSIL Hackathon 2026", size=48, color=WHITE, bold=True)

text_box(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.8),
         "Taiwan Hub  ·  全球健康創新黑客松", size=28, color=TEAL, bold=True)

text_box(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.6),
         "Building High-Value Health Systems: Leveraging AI", size=20, color=GOLD)

accent_line(slide, Inches(1.2), Inches(5.1), Inches(3), TEAL)

text_box(slide, Inches(1.2), Inches(5.5), Inches(10), Inches(0.5),
         "Harvard T.H. Chan School of Public Health  ×  National Taiwan University",
         size=16, color=MED_GRAY)

text_box(slide, Inches(1.2), Inches(6.1), Inches(10), Inches(0.5),
         "April 10–11, 2026  ·  NTU School of Public Health",
         size=16, color=MED_GRAY)

slide_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — What is HSIL Hackathon?
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Accent bar
rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CRIMSON)

text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "What is HSIL Hackathon?", size=36, color=CRIMSON, bold=True)
accent_line(slide, Inches(0.8), Inches(1.25), Inches(2.5), TEAL)

# Left column
items_left = [
    "Annual health innovation hackathon by Harvard HSIL",
    "Now in its 7th edition — 30+ cities worldwide",
    "2026 is Taiwan's FIRST time participating",
    "Co-organized with National Taiwan University (NTU)",
    "Theme: Leveraging AI for high-value health systems",
]
add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(6.5), Inches(4.5),
                items_left, size=20, color=DARK, bullet_color=TEAL)

# Right side — poster image
try:
    slide.shapes.add_picture(os.path.join(ASSETS, "hackathon-poster.png"),
                             Inches(8.0), Inches(1.2), height=Inches(5.5))
except:
    pass

slide_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — Quick Facts
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TEAL)

text_box(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "Quick Facts", size=36, color=CRIMSON, bold=True)
accent_line(slide, Inches(0.8), Inches(1.25), Inches(2), TEAL)

facts = [
    ("📅  Dates",          "April 10–11, 2026 (Fri & Sat)"),
    ("📍  Location",       "NTU School of Public Health"),
    ("⏱  Duration",       "2-Day Hackathon"),
    ("👥  Team Size",      "3–5 members per team"),
    ("🎤  Deliverable",    "3-min pitch + 2-min Q&A"),
    ("💰  Cost",           "Completely FREE (meals included)"),
    ("🗣  Language",       "Primarily Mandarin Chinese"),
    ("📝  Registration",   "Deadline: March 31, 2026"),
]

cols = 2
card_w = Inches(5.5)
card_h = Inches(1.15)
x_starts = [Inches(0.8), Inches(6.8)]

for i, (label, detail) in enumerate(facts):
    col = i % cols
    row = i // cols
    x = x_starts[col]
    y = Inches(1.7) + row * Inches(1.35)

    card = rounded_rect(slide, x, y, card_w, card_h, LIGHT_GRAY)

    text_box(slide, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.5), Inches(0.45),
             label, size=17, color=TEAL, bold=True)
    text_box(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.5), Inches(0.45),
             detail, size=17, color=DARK)

slide_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — Why Join? (Highlights)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GOLD)

text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
         "Why Join?", size=36, color=CRIMSON, bold=True)
accent_line(slide, Inches(0.8), Inches(1.25), Inches(1.5), GOLD)

highlights = [
    ("🆓  Completely Free",              "No registration fee — meals & materials provided"),
    ("💻  No Coding Required",           "All backgrounds welcome; deliverable is a pitch, not code"),
    ("🎓  Harvard Bootcamp",             "Top 2 teams get invited to HSIL Harvard Bootcamp"),
    ("🌍  Global Demo Day",              "Present at Harvard Demo Day on June 19, 2026"),
    ("🤝  Cross-Disciplinary",           "Network with students from medicine, CS, business & more"),
    ("🏥  Real Health Challenges",       "Solve actual healthcare problems with AI-driven solutions"),
]

card_w = Inches(5.8)
card_h = Inches(1.2)
x_starts = [Inches(0.8), Inches(6.9)]

for i, (title, desc) in enumerate(highlights):
    col = i % 2
    row = i // 2
    x = x_starts[col]
    y = Inches(1.7) + row * Inches(1.6)

    card = rounded_rect(slide, x, y, card_w, card_h, SOFT_TEAL_BG)

    text_box(slide, x + Inches(0.3), y + Inches(0.12), card_w - Inches(0.5), Inches(0.45),
             title, size=18, color=TEAL, bold=True)
    text_box(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.5), Inches(0.5),
             desc, size=16, color=DARK)

slide_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — Challenge Themes
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TEAL)

text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Challenge Themes", size=36, color=WHITE, bold=True)
text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "運用 AI 打造高價值醫療體系  —  Building High-Value Health Systems: Leveraging AI",
         size=17, color=TEAL)

themes = [
    "📊  EHR Data Analysis",
    "🤖  AI Diagnostic Tools",
    "💬  Clinical Chatbots",
    "👨‍⚕️  Medical Workforce Optimization",
    "❤️  Cardiovascular Disease",
    "🎗  Cancer Care",
    "🧠  Mental Health Services",
    "🛡  Disease Prevention & Surveillance",
]

card_w = Inches(5.5)
card_h = Inches(0.85)
x_starts = [Inches(0.8), Inches(6.9)]

for i, theme in enumerate(themes):
    col = i % 2
    row = i // 2
    x = x_starts[col]
    y = Inches(1.9) + row * Inches(1.15)

    card = rounded_rect(slide, x, y, card_w, card_h, RGBColor(0x25, 0x2B, 0x3B))
    text_box(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.5), Inches(0.5),
             theme, size=19, color=WHITE, bold=False)

slide_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — Schedule / Timeline
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), CRIMSON)

text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
         "Event Timeline", size=36, color=CRIMSON, bold=True)
accent_line(slide, Inches(0.8), Inches(1.25), Inches(2), CRIMSON)

timeline = [
    ("Now → Mar 31",   "Registration Open",                      "Online via Airtable form",                   TEAL),
    ("Early April",    "Acceptance & Pre-Event",                  "Team confirmation + preparation materials",   GOLD),
    ("Apr 10 (Day 1)", "Opening · Team Formation · Brainstorm",   "8:00 AM – 5:30 PM at NTU",                  CRIMSON),
    ("Apr 11 (Day 2)", "Solution Refinement · Pitch · Awards",    "8:30 AM – 5:00 PM at NTU",                  CRIMSON),
    ("May–June",       "Post-Event: Refinement & Harvard Demo",   "June 19 Demo Day at Harvard",                TEAL),
]

for i, (date, title, detail, dot_color) in enumerate(timeline):
    y = Inches(1.7) + i * Inches(1.05)

    # Timeline dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), y + Inches(0.12), Inches(0.28), Inches(0.28))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()

    # Vertical line (except last)
    if i < len(timeline) - 1:
        line = rect(slide, Inches(1.89), y + Inches(0.4), Inches(0.04), Inches(0.77), LIGHT_GRAY)

    # Date
    text_box(slide, Inches(0.3), y + Inches(0.08), Inches(1.4), Inches(0.4),
             date, size=14, color=MED_GRAY, bold=True, align=PP_ALIGN.RIGHT)

    # Title & detail
    text_box(slide, Inches(2.4), y, Inches(6), Inches(0.38),
             title, size=18, color=DARK, bold=True)
    text_box(slide, Inches(2.4), y + Inches(0.38), Inches(6), Inches(0.35),
             detail, size=14, color=MED_GRAY)

slide_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — Judging Criteria
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GOLD)

text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
         "Judging Criteria", size=36, color=CRIMSON, bold=True)
accent_line(slide, Inches(0.8), Inches(1.25), Inches(2), GOLD)

text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.4),
         "Each criterion scored 1–5 pts  ·  Total: 30 points", size=16, color=MED_GRAY)

criteria = [
    ("Innovation 創新性",        "Novel and creative approach to the problem"),
    ("Feasibility 可行性",       "Realistic implementation within constraints"),
    ("Health Impact 健康影響力",  "Meaningful contribution to health outcomes"),
    ("Scalability 可擴展性",     "Potential to scale across populations"),
    ("Presentation 簡報品質",     "Clear, compelling 3-minute pitch delivery"),
    ("AI Integration AI應用",    "Effective use of AI technologies"),
]

card_w = Inches(3.7)
card_h = Inches(1.8)
cols = 3

for i, (name, desc) in enumerate(criteria):
    col = i % cols
    row = i // cols
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(2.2) + row * Inches(2.2)

    card = rounded_rect(slide, x, y, card_w, card_h, SOFT_CRIM_BG)

    # Score badge
    badge = rounded_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.45), CRIMSON)
    text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.45),
             "5 pts", size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    text_box(slide, x + Inches(1.15), y + Inches(0.2), card_w - Inches(1.4), Inches(0.45),
             name, size=15, color=CRIMSON, bold=True)
    text_box(slide, x + Inches(0.2), y + Inches(0.85), card_w - Inches(0.4), Inches(0.8),
             desc, size=14, color=DARK)

slide_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack (Our Website)
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TEAL)

text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         "Our Website — Tech Stack", size=36, color=WHITE, bold=True)
text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.4),
         "Built with modern web technologies to deliver a polished event experience",
         size=16, color=MED_GRAY)

stacks = [
    ("⚛️  Framework",    "React 18 + TypeScript\nVite 5 (blazing-fast builds)"),
    ("🎨  Styling",      "Tailwind CSS 3\nshadcn/ui + Radix UI"),
    ("✨  Animations",   "Framer Motion\nScroll-triggered reveals"),
    ("📱  Responsive",   "Mobile-first design\nAdaptive navigation"),
    ("🔧  Dev Tools",    "ESLint · Vitest\nReact Testing Library"),
    ("🌐  Features",     "Live countdown timer\nBilingual (中文 / English)"),
]

card_w = Inches(3.7)
card_h = Inches(2.0)

for i, (title, detail) in enumerate(stacks):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.8) + row * Inches(2.5)

    card = rounded_rect(slide, x, y, card_w, card_h, RGBColor(0x25, 0x2B, 0x3B))

    text_box(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.5), Inches(0.45),
             title, size=18, color=TEAL, bold=True)

    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.75), card_w - Inches(0.5), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(detail.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(4)

slide_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — Website Demo / Key Sections
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), TEAL)

text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
         "Website Features", size=36, color=CRIMSON, bold=True)
accent_line(slide, Inches(0.8), Inches(1.25), Inches(2.5), TEAL)

sections = [
    ("Hero Section",       "Animated entrance with live countdown\ntimer to registration deadline"),
    ("About & Quick Facts","Event overview, 6-card fact grid,\n5-step process timeline"),
    ("Themes",             "8 AI-health challenge categories\nwith dark-mode card design"),
    ("Schedule",           "Color-coded vertical timeline\nfrom registration to Harvard Demo"),
    ("Judging & FAQ",      "6 scored criteria + 8 accordion\nFAQ items for participants"),
    ("CTA & Footer",       "Prominent register button with glow\neffect + social media links"),
]

card_w = Inches(3.7)
card_h = Inches(1.7)

for i, (title, detail) in enumerate(sections):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.7) + row * Inches(2.2)

    card = rounded_rect(slide, x, y, card_w, card_h, LIGHT_GRAY)

    text_box(slide, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.5), Inches(0.45),
             title, size=17, color=TEAL, bold=True)

    txBox = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.65), card_w - Inches(0.5), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(detail.split('\n')):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_before = Pt(3)

slide_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════════
# SLIDE 10 — Thank You / CTA
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

# Decorative
rect(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_H, CRIMSON)
rect(slide, Inches(0.4), Inches(0), Inches(0.15), SLIDE_H, TEAL)

text_box(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.2),
         "準備好改變世界了嗎？", size=44, color=WHITE, bold=True)
text_box(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.8),
         "Ready to Change the World?", size=28, color=TEAL, bold=True)

accent_line(slide, Inches(1.2), Inches(3.7), Inches(3), GOLD)

text_box(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.5),
         "📝  Register by March 31, 2026", size=22, color=WHITE)
text_box(slide, Inches(1.2), Inches(4.8), Inches(10), Inches(0.5),
         "📍  NTU School of Public Health  ·  April 10–11", size=18, color=MED_GRAY)
text_box(slide, Inches(1.2), Inches(5.3), Inches(10), Inches(0.5),
         "📸  Instagram: @hsilhackathon.taiwan", size=18, color=MED_GRAY)

# Logos at bottom
try:
    slide.shapes.add_picture(os.path.join(ASSETS, "hsil-logo.png"),
                             Inches(1.2), Inches(6.2), height=Inches(0.7))
except:
    pass

text_box(slide, Inches(4), Inches(6.7), Inches(8), Inches(0.4),
         "Thank you!  ·  謝謝！", size=20, color=MED_GRAY, align=PP_ALIGN.LEFT)

slide_number(slide, 10, TOTAL_SLIDES)

# ── Save ────────────────────────────────────────────────
output_path = "/home/user/HSIL-Website/HSIL_Hackathon_2026_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {TOTAL_SLIDES}")
